import os
import json
import hashlib
from typing import List, Dict, Any, Optional
import pandas as pd
import chromadb
from chromadb.config import Settings
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import PyPDFLoader, TextLoader
from langchain_openai import OpenAIEmbeddings
from docx import Document
import streamlit as st


class RAGManager:
    _chroma_client = None
    _chroma_settings = None
    _chroma_dir = None

    def __init__(self, data_dir: str = "data"):
        self.data_dir = data_dir
        self.chroma_dir = os.path.join(data_dir, "chroma_db")
        os.makedirs(self.chroma_dir, exist_ok=True)

        # テキスト分割器の初期化（最初に実行）
        self.text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)

        # 埋め込みモデルの初期化
        try:
            self.embeddings = OpenAIEmbeddings()
        except Exception as e:
            st.warning("OpenAI API key not set. Using default embeddings.")
            self.embeddings = None

        # ChromaDBクライアントの初期化（シングルトンパターン）
        self.client = None
        self.chroma_available = False

        # クラス変数として同じ設定を共有
        if RAGManager._chroma_dir != self.chroma_dir or RAGManager._chroma_client is None:
            self._initialize_chroma_client()
        else:
            # 既存のクライアントを再利用
            self.client = RAGManager._chroma_client
            self.chroma_available = self.client is not None

    def _initialize_chroma_client(self):
        """ChromaDBクライアントの初期化"""
        try:
            # 既存のクライアントがある場合はリセット
            if RAGManager._chroma_client is not None:
                try:
                    RAGManager._chroma_client = None
                except Exception:
                    pass

            # ディレクトリの権限を確実に設定
            try:
                os.chmod(self.chroma_dir, 0o755)
                for root, dirs, files in os.walk(self.chroma_dir):
                    for d in dirs:
                        os.chmod(os.path.join(root, d), 0o755)
                    for f in files:
                        os.chmod(os.path.join(root, f), 0o644)
            except Exception:
                pass  # 権限設定に失敗しても続行

            # ChromaDB設定（readonly対策）
            settings = Settings(
                allow_reset=True,
                anonymized_telemetry=False
            )

            RAGManager._chroma_settings = settings
            RAGManager._chroma_dir = self.chroma_dir
            RAGManager._chroma_client = chromadb.PersistentClient(path=self.chroma_dir, settings=settings)

            self.client = RAGManager._chroma_client
            self.chroma_available = True
        except Exception as chroma_error:
            error_msg = str(chroma_error)

            # 各種エラーに対応した自動復旧
            if ("readonly database" in error_msg or
                "no such table: databases" in error_msg or
                "database is locked" in error_msg or
                "already exists" in error_msg or
                "different settings" in error_msg):
                st.warning("🔄 ChromaDBデータベースを再初期化しています...")
                import shutil
                if os.path.exists(self.chroma_dir):
                    shutil.rmtree(self.chroma_dir)
                os.makedirs(self.chroma_dir, exist_ok=True)

                # ディレクトリの権限を確実に設定
                try:
                    os.chmod(self.chroma_dir, 0o755)
                except Exception:
                    pass

                try:
                    settings = Settings(
                        allow_reset=True,
                        anonymized_telemetry=False
                    )
                    self.client = chromadb.PersistentClient(path=self.chroma_dir, settings=settings)
                    self.chroma_available = True
                    st.success("✅ ChromaDB データベースが正常に再初期化されました")
                except Exception:
                    st.error("⚠️ ChromaDB再初期化に失敗: RAG機能は無効になります")
                    self.client = None
                    self.chroma_available = False
            else:
                st.error(f"⚠️ ChromaDB初期化エラー: {error_msg}")
                st.error("RAG機能は無効になります")
                self.client = None
                self.chroma_available = False



    def get_or_create_collection(self, product_name: str):
        if not self.chroma_available:
            st.error("❌ ChromaDB が利用できません")
            return None

        collection_name = f"product_{product_name.lower().replace(' ', '_')}"
        try:
            collection = self.client.get_collection(name=collection_name)
            if st.secrets.get("DEBUG_MODE", False):
                st.success(f"✅ 既存コレクション取得: {collection_name}")
        except Exception as get_error:
            get_error_msg = str(get_error)

            # readonly エラーの場合は再初期化を試行
            if "readonly database" in get_error_msg or "database is locked" in get_error_msg:
                st.warning("🔄 データベースアクセスエラー、再初期化を試行...")
                self._reinitialize_chroma()
                if not self.chroma_available:
                    return None

            try:
                collection = self.client.create_collection(name=collection_name, metadata={"hnsw:space": "cosine"})
                if st.secrets.get("DEBUG_MODE", False):
                    st.success(f"✅ 新規コレクション作成: {collection_name}")
            except Exception as create_error:
                create_error_msg = str(create_error)

                # readonly エラーの場合は再初期化を試行
                if "readonly database" in create_error_msg or "database is locked" in create_error_msg:
                    st.warning("🔄 コレクション作成でアクセスエラー、再初期化を試行...")
                    self._reinitialize_chroma()
                    if not self.chroma_available:
                        return None

                    # 再初期化後に再試行
                    try:
                        collection = self.client.create_collection(name=collection_name, metadata={"hnsw:space": "cosine"})
                        if st.secrets.get("DEBUG_MODE", False):
                            st.success(f"✅ 再初期化後コレクション作成成功: {collection_name}")
                    except Exception:
                        st.error(f"❌ 再初期化後もコレクション作成失敗: {collection_name}")
                        return None
                else:
                    st.error(f"❌ コレクション作成失敗: {collection_name}")
                    st.error(f"取得エラー: {get_error}")
                    st.error(f"作成エラー: {create_error}")
                    return None

        return collection

    def _reinitialize_chroma(self):
        """ChromaDBの再初期化"""
        try:
            # 既存のクライアントをクリア
            RAGManager._chroma_client = None

            import shutil
            if os.path.exists(self.chroma_dir):
                shutil.rmtree(self.chroma_dir)
            os.makedirs(self.chroma_dir, exist_ok=True)

            # ディレクトリの権限を設定
            try:
                os.chmod(self.chroma_dir, 0o755)
            except Exception:
                pass

            # 新しいクライアントを作成
            settings = Settings(
                allow_reset=True,
                anonymized_telemetry=False
            )

            RAGManager._chroma_settings = settings
            RAGManager._chroma_dir = self.chroma_dir
            RAGManager._chroma_client = chromadb.PersistentClient(path=self.chroma_dir, settings=settings)

            self.client = RAGManager._chroma_client
            self.chroma_available = True
            st.info("✅ ChromaDB再初期化完了")
        except Exception as e:
            st.error(f"❌ ChromaDB再初期化失敗: {e}")
            RAGManager._chroma_client = None
            self.client = None
            self.chroma_available = False

    def get_file_hash(self, file_content: bytes) -> str:
        return hashlib.md5(file_content).hexdigest()

    def _is_duplicate_file(self, product_name: str, file_hash: str) -> bool:
        """ファイルハッシュによる重複チェック"""
        if not self.chroma_available:
            return False

        try:
            collection = self.get_or_create_collection(product_name)
            if collection is None:
                return False

            results = collection.get(where={"file_hash": file_hash})
            return bool(results and results["ids"])
        except Exception:
            return False

    def extract_text_from_file(self, file_path: str, file_type: str) -> str:
        try:
            if file_type == "txt":
                with open(file_path, "r", encoding="utf-8") as f:
                    return f.read()
            elif file_type == "pdf":
                loader = PyPDFLoader(file_path)
                documents = loader.load()
                return "\n".join([doc.page_content for doc in documents])
            elif file_type == "docx":
                doc = Document(file_path)
                return "\n".join([paragraph.text for paragraph in doc.paragraphs])
            elif file_type == "pptx":
                from pptx import Presentation

                prs = Presentation(file_path)
                text_content = []

                for slide_num, slide in enumerate(prs.slides, 1):
                    slide_text = f"=== スライド {slide_num} ===\n"

                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text += shape.text.strip() + "\n"

                        # 表がある場合の処理
                        if hasattr(shape, "table"):
                            table = shape.table
                            for row in table.rows:
                                row_text = " | ".join([cell.text.strip() for cell in row.cells])
                                slide_text += row_text + "\n"

                    text_content.append(slide_text)

                return "\n\n".join(text_content)
            elif file_type == "html":
                with open(file_path, "r", encoding="utf-8") as f:
                    content = f.read()
                from bs4 import BeautifulSoup

                soup = BeautifulSoup(content, "html.parser")
                return soup.get_text()
            elif file_type == "csv":
                df = pd.read_csv(file_path)

                # CSVの各行をQ&A形式として処理
                qa_pairs = []

                # 列数をチェック
                if len(df.columns) >= 2:
                    # 2列以上の場合、最初の列を質問、2番目の列を回答として処理
                    # 3列目は参照データ（RAG処理には含めない）
                    question_col = df.columns[0]
                    answer_col = df.columns[1]

                    for index, row in df.iterrows():
                        question = str(row[question_col]).strip()
                        answer = str(row[answer_col]).strip()

                        if question and answer and question != "nan" and answer != "nan":
                            # 参照データはRAG処理に含めない
                            qa_text = f"質問: {question}\n回答: {answer}"
                            qa_pairs.append(qa_text)
                else:
                    # 1列の場合は従来通りの処理
                    return df.to_string()

                return "\n\n".join(qa_pairs) if qa_pairs else df.to_string()
            else:
                return ""
        except Exception as e:
            st.error(f"Error reading file: {e}")
            return ""

    def add_csv_as_qa_pairs(self, product_name: str, file_name: str, file_content: bytes) -> bool:
        """CSVファイルの各行を個別のQ&Aペアとして処理"""
        if not self.chroma_available:
            st.error("❌ ChromaDB が利用できないため、ファイルを追加できません")
            st.error("💡 アプリを再起動してください")
            return False

        try:
            collection = self.get_or_create_collection(product_name)
            if collection is None:
                st.error(f"❌ 商材「{product_name}」のコレクション作成に失敗しました")
                return False

            file_hash = self.get_file_hash(file_content)

            # 重複チェック
            if self._is_duplicate_file(product_name, file_hash):
                st.warning(f"⚠️ CSVファイル「{file_name}」は既に追加されています（重複スキップ）")
                return True

            temp_file_path = os.path.join(self.data_dir, f"temp_{file_hash}.csv")
            with open(temp_file_path, "wb") as f:
                f.write(file_content)

            df = pd.read_csv(temp_file_path)

            if len(df.columns) >= 2:
                # カラム名を検証して適切に設定（類似語も検出）
                def find_column_by_keywords(columns, keywords):
                    """キーワードリストからカラムを検索"""
                    for col in columns:
                        if any(keyword in str(col).lower() for keyword in keywords):
                            return col
                    return None

                # カラム候補の定義
                question_keywords = ["質問", "question", "q", "問", "問い"]
                answer_keywords = ["回答", "答え", "answer", "a", "応答", "回"]
                reference_keywords = ["参考", "参照", "reference", "ref", "リンク", "link", "url", "出典"]

                # 各カラムを検索
                question_col = find_column_by_keywords(df.columns, question_keywords)
                answer_col = find_column_by_keywords(df.columns, answer_keywords)
                reference_col = find_column_by_keywords(df.columns, reference_keywords)

                # 見つからない場合は位置ベースでフォールバック
                if not question_col:
                    question_col = df.columns[0]
                if not answer_col:
                    answer_col = df.columns[1]
                if not reference_col and len(df.columns) >= 3:
                    reference_col = df.columns[2]

                # デバッグ情報表示
                if st.secrets.get("DEBUG_MODE", False):
                    st.write(f"🔍 **CSV カラム自動検出結果**:")
                    st.write(f"- 検出されたカラム: {list(df.columns)}")
                    st.write(f"- 質問カラム: '{question_col}'")
                    st.write(f"- 回答カラム: '{answer_col}'")
                    st.write(f"- 参照カラム: '{reference_col}'")

                qa_count = 0
                for index, row in df.iterrows():
                    question = str(row[question_col]).strip()
                    answer = str(row[answer_col]).strip()

                    if question and answer and question != "nan" and answer != "nan":
                        # RAG処理用テキスト（参照データは含めない）
                        qa_text = f"質問: {question}\n回答: {answer}"
                        doc_id = f"{file_hash}_qa_{index}"

                        # 参照データの取得（あれば）
                        reference_data = ""
                        if reference_col and pd.notna(row[reference_col]):
                            reference_data = str(row[reference_col]).strip()

                        # メタデータに参照データを保存（検索結果表示用）
                        metadata = {
                            "file_name": file_name,
                            "file_hash": file_hash,
                            "qa_index": index,
                            "product": product_name,
                            "question": question,
                            "answer": answer,
                            "type": "qa_pair",
                        }

                        # 参照データがある場合のみ追加
                        if reference_data:
                            metadata["reference"] = reference_data

                        try:
                            collection.add(
                                documents=[qa_text],  # 参照データはRAG処理に含めない
                                metadatas=[metadata],
                                ids=[doc_id],
                            )
                        except Exception as add_error:
                            if "readonly database" in str(add_error) or "database is locked" in str(add_error):
                                st.warning("🔄 データベースエラー、再初期化して再試行...")
                                self._reinitialize_chroma()
                                collection = self.get_or_create_collection(product_name)
                                if collection:
                                    collection.add(
                                        documents=[qa_text],
                                        metadatas=[metadata],
                                        ids=[doc_id],
                                    )
                            else:
                                raise add_error
                        qa_count += 1

                os.remove(temp_file_path)
                # 参照データ付きの件数も表示
                reference_count = sum(1 for _, row in df.iterrows()
                                    if reference_col and pd.notna(row[reference_col]) and str(row[reference_col]).strip())

                # 詳細なサマリー情報を表示
                if st.secrets.get("DEBUG_MODE", False):
                    st.write(f"🔍 **CSV処理サマリー**: {file_name}")
                    st.write(f"- 検出カラム: {list(df.columns)}")
                    st.write(f"- 使用カラム: 質問='{question_col}', 回答='{answer_col}', 参照='{reference_col}'")
                    st.write(f"- 総行数: {len(df)}, 処理済み: {qa_count}")

                if reference_count > 0:
                    st.success(f"✅ {qa_count}組のQ&Aペア（うち{reference_count}件は参照データ付き）を追加しました")
                else:
                    st.success(f"✅ {qa_count}組のQ&Aペアを追加しました")
                return True
            else:
                st.error("CSVファイルには質問と回答の2列が必要です")
                os.remove(temp_file_path)
                return False

        except Exception as e:
            # エラー時も一時ファイルをクリーンアップ
            if 'temp_file_path' in locals() and os.path.exists(temp_file_path):
                os.remove(temp_file_path)
            st.error(f"❌ CSV処理エラー: {str(e)}")
            return False

    def add_document(self, product_name: str, file_name: str, file_content: bytes, file_type: str) -> bool:
        if not self.chroma_available:
            st.error("❌ ChromaDB が利用できないため、ファイルを追加できません")
            st.error("💡 アプリを再起動してください")
            return False

        try:
            collection = self.get_or_create_collection(product_name)
            if collection is None:
                st.error(f"❌ 商材「{product_name}」のコレクション作成に失敗しました")
                return False

            file_hash = self.get_file_hash(file_content)

            # 重複チェック
            if self._is_duplicate_file(product_name, file_hash):
                st.warning(f"⚠️ ファイル「{file_name}」は既に追加されています（重複スキップ）")
                return True

            temp_file_path = os.path.join(self.data_dir, f"temp_{file_hash}.{file_type}")
            with open(temp_file_path, "wb") as f:
                f.write(file_content)

            text_content = self.extract_text_from_file(temp_file_path, file_type)

            if not text_content or not text_content.strip():
                st.error(f"❌ ファイル「{file_name}」からテキストを抽出できませんでした")
                os.remove(temp_file_path)
                return False

            if text_content:
                chunks = self.text_splitter.split_text(text_content)

                if st.secrets.get("DEBUG_MODE", False):
                    st.info(f"📄 ファイル処理: {file_name}")
                    st.info(f"📊 テキスト長: {len(text_content)}文字")
                    st.info(f"🧩 チャンク数: {len(chunks)}個")

                for i, chunk in enumerate(chunks):
                    doc_id = f"{file_hash}_{i}"
                    try:
                        collection.add(
                            documents=[chunk],
                            metadatas=[
                                {"file_name": file_name, "file_hash": file_hash, "chunk_index": i, "product": product_name}
                            ],
                            ids=[doc_id],
                        )
                    except Exception as add_error:
                        if "readonly database" in str(add_error) or "database is locked" in str(add_error):
                            st.warning("🔄 データベースエラー、再初期化して再試行...")
                            self._reinitialize_chroma()
                            collection = self.get_or_create_collection(product_name)
                            if collection:
                                collection.add(
                                    documents=[chunk],
                                    metadatas=[
                                        {"file_name": file_name, "file_hash": file_hash, "chunk_index": i, "product": product_name}
                                    ],
                                    ids=[doc_id],
                                )
                        else:
                            raise add_error

                if st.secrets.get("DEBUG_MODE", False):
                    st.success(f"✅ {len(chunks)}個のチャンクをChromeDBに追加完了")

            os.remove(temp_file_path)
            return True

        except Exception as e:
            # エラー時も一時ファイルをクリーンアップ
            if 'temp_file_path' in locals() and os.path.exists(temp_file_path):
                os.remove(temp_file_path)
            st.error(f"❌ ドキュメント追加エラー: {e}")
            return False

    def remove_document(self, product_name: str, file_name: str) -> bool:
        if not self.chroma_available:
            st.error("❌ ChromaDB が利用できないため、ファイルを削除できません")
            return False

        try:
            collection = self.get_or_create_collection(product_name)
            if collection is None:
                return False

            results = collection.get(where={"file_name": file_name})

            if results and results["ids"]:
                try:
                    collection.delete(ids=results["ids"])
                    return True
                except Exception as delete_error:
                    if "readonly database" in str(delete_error) or "database is locked" in str(delete_error):
                        st.warning("🔄 データベースエラー、再初期化して再試行...")
                        self._reinitialize_chroma()
                        collection = self.get_or_create_collection(product_name)
                        if collection:
                            results = collection.get(where={"file_name": file_name})
                            if results and results["ids"]:
                                collection.delete(ids=results["ids"])
                                return True
                    else:
                        raise delete_error
            return False

        except Exception as e:
            st.error(f"Error removing document: {e}")
            return False

    def search(self, product_name: str, query: str, top_k: int = 5) -> List[Dict[str, Any]]:
        if not self.chroma_available:
            return []

        try:
            collection = self.get_or_create_collection(product_name)
            if collection is None:
                return []

            results = collection.query(query_texts=[query], n_results=top_k)

            search_results = []
            if results["documents"] and results["documents"][0]:
                for i in range(len(results["documents"][0])):
                    distance = results["distances"][0][i] if results["distances"] else 0.0

                    # ChromaDBは cosine_distance = 1.0 - cosine_similarity を返す
                    # したがって cosine_similarity = 1.0 - cosine_distance
                    cosine_similarity = 1.0 - distance

                    # 数値精度の問題で負になる可能性を考慮（実際は稀）
                    similarity_score = max(0.0, min(1.0, cosine_similarity))

                    search_results.append(
                        {
                            "content": results["documents"][0][i],
                            "metadata": results["metadatas"][0][i] if results["metadatas"] else {},
                            "distance": distance,
                            "similarity_score": similarity_score,
                        }
                    )

            # デバッグ情報を出力（開発時のみ）
            if st.secrets.get("DEBUG_MODE", False):
                st.write(f"🔍 検索クエリ: '{query}' (上位{len(search_results)}件)")
                st.write("**ChromaDB計算式**: cosine_similarity = 1.0 - cosine_distance")

                for i, result in enumerate(search_results[:3]):
                    distance = result['distance']
                    similarity = result['similarity_score']
                    cosine_sim_check = 1.0 - distance
                    st.write(f"結果{i+1}: 距離={distance:.4f}, 類似度={similarity:.4f} (検算: {cosine_sim_check:.4f})")
                    st.write(f"内容: {result['content'][:100]}...")

                # ソート前後の順序確認
                if len(search_results) > 1:
                    distances = [r['distance'] for r in search_results]
                    is_sorted = all(distances[i] <= distances[i+1] for i in range(len(distances)-1))
                    st.write(f"📊 ChromaDB結果ソート状況: {'✅ 既にソート済み' if is_sorted else '❌ ソートが必要'}")

                # 距離と類似度の範囲確認
                min_dist, max_dist = min(distances), max(distances)
                min_sim = min(r['similarity_score'] for r in search_results)
                max_sim = max(r['similarity_score'] for r in search_results)
                st.write(f"📊 範囲確認: 距離=[{min_dist:.4f}, {max_dist:.4f}], 類似度=[{min_sim:.4f}, {max_sim:.4f}]")

            # ChromaDBは既に距離順（小さい順）で返すが、念のため明示的にソート
            search_results.sort(key=lambda x: x["distance"])

            # 関連度の低い結果をフィルタリング
            similarity_threshold = 0.7  # デフォルト値
            try:
                # 現在のRAG設定から閾値を取得
                _, rag_config = get_current_rag_config()
                similarity_threshold = getattr(rag_config, 'similarity_threshold', 0.7)
            except Exception as e:
                # デバッグモード時のみエラー表示
                if st.secrets.get("DEBUG_MODE", False):
                    st.warning(f"RAG設定取得エラー（デフォルト値0.7を使用): {e}")
                similarity_threshold = 0.7

            # 類似度が閾値以上のものを保持（より直感的）
            filtered_results = [r for r in search_results if r["similarity_score"] >= similarity_threshold]

            if st.secrets.get("DEBUG_MODE", False) and len(filtered_results) < len(search_results):
                st.write(f"⚠️ 関連度の低い結果を {len(search_results) - len(filtered_results)} 件除外しました")

            # フィルタリング後の結果処理
            if filtered_results:
                return filtered_results
            elif search_results:
                # フィルター後に結果がない場合、最も関連度の高い1件を返す
                if st.secrets.get("DEBUG_MODE", False):
                    st.warning(f"全結果が閾値({similarity_threshold})を下回りました。最高スコア結果を返します。")
                return search_results[:1]
            else:
                return []

        except Exception as e:
            st.error(f"Error searching: {e}")
            return []

    def list_documents(self, product_name: str) -> List[str]:
        if not self.chroma_available:
            return []

        try:
            collection = self.get_or_create_collection(product_name)
            if collection is None:
                return []
            results = collection.get()

            if results and results["metadatas"]:
                file_names = set()
                for metadata in results["metadatas"]:
                    if "file_name" in metadata:
                        file_names.add(metadata["file_name"])
                return list(file_names)
            return []

        except Exception as e:
            st.error(f"Error listing documents: {e}")
            return []

    def list_products(self) -> List[str]:
        if not self.chroma_available:
            return []

        try:
            collections = self.client.list_collections()
            products = []
            for collection in collections:
                if collection.name.startswith("product_"):
                    product_name = collection.name.replace("product_", "").replace("_", " ").title()
                    products.append(product_name)
            return products
        except Exception as e:
            st.error(f"Error listing products: {e}")
            return []
